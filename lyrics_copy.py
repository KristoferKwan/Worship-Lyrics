from bs4 import BeautifulSoup
import urllib.request
import pptx
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_FILL
import time
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys

WHITE_COLOR = RGBColor(255, 255, 255)
BLACK_COLOR = RGBColor(0, 0, 0)
COLOR = BLACK_COLOR

def remove_markup(line):
    line = line.replace('&quot;', '')
    line = line.replace('<br>', '')
    line = line.replace("verse: ", '')
    line = line.replace("\r", "\n")

    while (line.find('<') != -1 and line.find('>') != -1):
        firstIndex = line.find("<")
        lastIndex = line.rfind(">")
        retLine = line[0:firstIndex] + line[lastIndex:-1]
        line = retLine

    # remove extraneous markoup
    line = line.replace('<', '')

    line = line.replace(">", "") 

    return line

def fix_typos(line):
    line = line.replace(" i ", " I ")
    line = line.replace('.', '')
    line = line.replace(';', '')
    return line

def capitalize_God(line):
    line = line.replace("you", "You")
    return line

'''
Return the url to the lyric page
'''
def search_azlyrics(search):
    search = search.replace(' ', '+')
    url = "http://search.azlyrics.com/search.php?q=" + search
    html  =  urllib.request.urlopen(url).read()
    data = html
    #print(data)
    soup = BeautifulSoup(data, "html.parser")    
    group = soup.findAll('table', class_='table-condensed')
    listSongs = list()
    count=0
    for links in group:
        for song in links.findAll('a'):
            songURL = song['href']
            if songURL[len(songURL)-1] == 'l':
                listSongs.append(song['href'])
    print(listSongs)
    if len(listSongs) > 0:
        return listSongs      
    return None
    
def azLyricsScraper(azLyricsUrl):
    '''
    Manually scrape texts
    by tag:
    <!-- Usage of azlyrics.com content by any third-party lyrics provider is prohibited by our licensing agreement. Sorry about that. -->

    return a list of lyric lines
    '''
    req = urllib.request.Request(azLyricsUrl)
    html  =  urllib.request.urlopen(req).read()
    #print(type(html))
    try:
        data = html.decode("utf-8", errors='replace')
    except Exception as ex:
        print("cant decode char")
    # We're not a lyrics provider!
    keytag = '<!-- Usage of azlyrics.com content by any third-party lyrics provider is prohibited by our licensing agreement. Sorry about that. -->'
    endtag = '</div>'
    lyrics = []
    add = False
    data = data.split('\n')

    for line in data:
        try:
            if add == True and endtag in line:
                break        
            if add == True:
                line = remove_markup(line)
                line = fix_typos(line)
                line = capitalize_God(line)
                lyrics.append(line)  
                #print(line)      
            if keytag in line:
                add = True
        except Exception as ex:
            print(ex)
           
    return lyrics 


def searchMetrolyrics(search):
    '''
    metro lyrics data
    '''
    search = search.replace(' ', '+')
    url = 'http://www.metrolyrics.com/search.html?search=' + search
    html  =  urlopen(url)
    data = html.read() 
    soup = BeautifulSoup(data, "html.parser")  
    group = soup.findAll(class_='content clearfix')
    
    for item in group:
        for song in item.findAll('a'):
            print (item)
            print()
        
    return 'http://www.metrolyrics.com/mighty-to-save-lyrics-hillsong.html'

def metrolyricsScraper(metrolyricsUrl):
    html  =  urlopen(metrolyricsUrl)
    data = html.read()    
    data = data.decode().split('<div id="lyrics-body-text" class="js-lyric-text"')
    song = data[1]
    verses = song.decode().split("<p class='verse'>")
    add = True
    song = []
    
    for verse in verses:
        if '</div>' in verse:
            verse = verse.decode().split('</div>')
            verse = verse[0].replace('<br>', '').replace('</p>', '')
            song.append(verse)
            print (verse)
        elif add:
            verse = verse.replace('<br>','').replace('</p>', '')
            song.append(verse)
            print (verse)
            print()
            
def get_lyrics(song):
    '''
    get lyrics from a azlyrics
    returns a string of the paragraphs
    '''
    # find the song page url
    azLyricsUrlList = search_azlyrics(song)
    numberSong = int(input("\n\nEnter the index of the desired song (within python list) where 0 is the first element: "))
    print(azLyricsUrlList[numberSong])
    azLyricsUrl = azLyricsUrlList[numberSong]
    # get list of lyrics from azlyrics.com
    lyrics = azLyricsScraper(azLyricsUrl)
    print(lyrics)
    #list of song paragraphs
    song = []
     
    # verses are separated by blank strings
    # append lines until reach an end of verse
    # marker
    currentVerse = ''
    for line in lyrics:
        line = line.replace("\n", "").strip()
        if line != '':
            # TRY
            # line = remove_markup(line)
            currentVerse += line + '\n'
        # reached the end of a verse
        elif line == "":
            print ('verse: ' +currentVerse)
            song.append(currentVerse)
            currentVerse = ''
    if currentVerse != '':
        song.append(currentVerse)


    return song            

def get_song_text(song, text):
    '''
    get song text from lyric site and
    add its paragraphs to a list
    '''
    try:
        text += get_lyrics(song)
        text.append(' ')
    except ValueError as e:
        print (e)     
    
def add_slide(prs, text):
    '''
    add one slide to the powerpoint
    '''
    # text = remove_past(text, '[')
    #text slide
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame

    p = tf.add_paragraph()
    p.text = text
    p.font.color.rgb = COLOR
    p.alignment = PP_ALIGN.CENTER

    # width = height = Inches(1.8)
    # top = Inches(1.5)
    # left = Inches(4)
    # txBox = slide.shapes.add_textbox(left, top, width, height)
    # tf = txBox.text_frame
    # print("\nThe following is the textframe settings:",tf.auto_size, tf.word_wrap)   
    # p = tf.add_paragraph()
    # p.text = text
    # p.font.color.rgb = COLOR
    # p.alignment = PP_ALIGN.CENTER
    
    # print (text)
    # print ('------------------')
    # p.font.size = Pt(30) 
    # p.font.bold = False   
    # p.font.name = 'Helvetica'
    
def make_power_point(text, title):
    '''
    makes a full powerpoint out of a list 
    of paragraphs
    '''
    prs = pptx.Presentation()
    #title slide
    title_slide_layout = prs.slide_layouts[0]
    #slide = prs.slides.add_slide(title_slide_layout)
    left = top = Inches(0)
    #slide.shapes.add_picture('C:/Users/ravicz/Desktop/Python/scraper/cruIntro.jpg', left, top)
    
    # Try
    for verse in text:
        if(verse != ''):
            add_slide(prs, verse)
           
    title_slide_layout = prs.slide_layouts[0]
    #slide = prs.slides.add_slide(title_slide_layout)
    left = top = Inches(0)
    #slide.shapes.add_picture('C:/Users/ravicz/Desktop/Python/scraper/cruIntro.jpg', left, top) 
    print("Save?\n")
    answer = input()
    if (answer in {'y', 'yes'}):
        prs.save(title) 

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("More args")
        sys.exit(1)
    elif sys.argv[1] == '-w':
        COLOR = WHITE_COLOR
        songs = " ".join(sys.argv[2:])
    else:
        songs = " ".join(sys.argv[1:])
    print(songs)
    songs = songs.split(',')
    print (songs)
    text = []

    for song in songs:
        get_song_text(song, text)

    songList = "-".join(songs)
    date = time.strftime('%m-%d').split('-')
    title = songList + '_' + date[0] + '_' + date[1] + '.pptx'
    make_power_point(text, title)